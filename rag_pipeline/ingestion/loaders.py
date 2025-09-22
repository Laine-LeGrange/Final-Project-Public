# Import standard libraries
from __future__ import annotations
from pathlib import Path
from typing import Tuple, List
import os
import io
import tempfile


# Import optional dependencies with graceful degradation

# PDF
try:
    from langchain_community.document_loaders import PyMuPDFLoader
    _HAVE_PYMUPDF = True
except Exception:
    PyMuPDFLoader = None
    _HAVE_PYMUPDF = False

# Word
try:
    from langchain_community.document_loaders import Docx2txtLoader
    _HAVE_DOCX2TXT = True
except Exception:
    Docx2txtLoader = None
    _HAVE_DOCX2TXT = False

# PowerPoint
try:
    from langchain_community.document_loaders import UnstructuredPowerPointLoader
    _HAVE_UNSTRUCTURED_PPT = True
except Exception:
    UnstructuredPowerPointLoader = None
    _HAVE_UNSTRUCTURED_PPT = False

# PPTX (python-pptx)
try:
    import pptx
    _HAVE_PYTHON_PPTX = True
except Exception:
    pptx = None
    _HAVE_PYTHON_PPTX = False

# OCR
try:
    from ..providers.ocr import image_to_text as ocr_image_to_text
    _HAVE_OCR_PROVIDER = True
except Exception:
    ocr_image_to_text = None
    _HAVE_OCR_PROVIDER = False

# ASR
try:
    from ..providers.asr import audio_to_text as whisper_audio_to_text
    _HAVE_ASR = True
except Exception:
    whisper_audio_to_text = None
    _HAVE_ASR = False

# Import config
try:
    from ..settings import load_config
    _HAVE_CONFIG = True
except Exception:
    load_config = lambda: {}
    _HAVE_CONFIG = False

# Pydub for audio extraction from video
try:
    from pydub import AudioSegment
    _HAVE_PYDUB = True
except Exception:
    AudioSegment = None
    _HAVE_PYDUB = False


# ---------------- SUPPORTED FORMATS ---------------- #
SUPPORTED_DOC   = {".txt", ".md"}
SUPPORTED_IMAGE = {".png", ".jpg", ".jpeg"}
SUPPORTED_AUDIO = {".mp3", ".wav", ".m4a", ".mp4", ".mpeg", ".mpga", ".webm"}
SUPPORTED_VIDEO = {".mp4", ".mov", ".mkv", ".webm"}
SUPPORTED_PRES  = {".ppt", ".pptx"}
SUPPORTED_PDF   = {".pdf"}
SUPPORTED_WORD  = {".doc", ".docx"}


# Get file extensions
def get_file_suffix(name: str) -> str:
    """ Return the file suffix (extension) in lowercase, or empty string if none. """
    return Path(name).suffix.lower()

# Write bytes to tempfile
def write_temp_file(data: bytes, suffix: str) -> str:
    """Write bytes to a temp file and return the path."""
    ntf = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    ntf.write(data)
    ntf.flush()
    ntf.close()
    return ntf.name

# ---------------- Loaders for specific formats ---------------- #

# PDF loader
def load_text_from_pdf(file_name: str, data: bytes) -> List[str]:
    """Load text from PDF using PyMuPDFLoader."""
    if not _HAVE_PYMUPDF:
        raise RuntimeError("PyMuPDFLoader not installed. Run: pip install langchain-community pymupdf")
    tmp = write_temp_file(data, suffix=".pdf")
    docs = PyMuPDFLoader(tmp).load()
    return [d.page_content for d in docs if (d.page_content or "").strip()]


# Word loader for .docx
def load_text_from_docx(file_name: str, data: bytes) -> List[str]:
    """Load text from DOCX using Docx2txtLoader."""
    if not _HAVE_DOCX2TXT:
        raise RuntimeError("Docx2txtLoader not installed. Run: pip install langchain-community docx2txt")
    tmp = write_temp_file(data, suffix=".docx")
    docs = Docx2txtLoader(tmp).load()
    return [d.page_content for d in docs if (d.page_content or "").strip()]


# PowerPoint loader for .ppt and .pptx
def load_text_from_ppt(file_name: str, data: bytes) -> List[str]:
    """
    Use UnstructuredPowerPointLoader (works for .ppt/.pptx),
    and fallback to python-pptx for .pptx files only.
    """
    ext = get_file_suffix(file_name)
    tmp_suffix = ".pptx" if ext == ".pptx" else ".ppt"
    tmp = write_temp_file(data, suffix=tmp_suffix)

    # Try UnstructuredPowerPointLoader first
    if _HAVE_UNSTRUCTURED_PPT:
        try:
            docs = UnstructuredPowerPointLoader(tmp).load()
            out = [d.page_content for d in docs if (d.page_content or "").strip()]
            if out:
                return out
        except Exception:
            pass

        # Fallback for .pptx using python-pptx
    if _HAVE_PYTHON_PPTX and ext == ".pptx":
        try:
            # Load presentation
            prs = pptx.Presentation(tmp)
            chunks: List[str] = []
            for i, slide in enumerate(prs.slides, start=1):
                lines: List[str] = []

                # Extract text from shapes
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text:
                            t = shape.text.strip()
                            if t:
                                lines.append(t)
                        elif hasattr(shape, "text_frame") and shape.text_frame:
                            
                            t = shape.text_frame.text.strip()
                            if t:
                                lines.append(t)
                    except Exception:
                        continue
                try:
                    # Extract notes if present
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame
                        if notes:
                            nt = notes.text.strip()
                            if nt:
                                lines.append(f"Notes: {nt}")
                except Exception:
                    pass
                # Only add slide if it has content
                if lines:
                    chunks.append(f"Slide {i}\n" + "\n".join(lines))
            return chunks or [f"[Empty presentation content] ({file_name})"]
        except Exception:
            pass

    # If all else fails
    return [f"[Presentation parse error or unsupported format] ({file_name})"]


# OCR for images
def load_image_text(file_name: str, data: bytes) -> List[str]:
    """ Use OCR to extract text from images. """
    if not _HAVE_OCR_PROVIDER or ocr_image_to_text is None:
        return [f"[OCR provider unavailable] ({file_name})"]
    text = ocr_image_to_text(image_bytes=data)
    text = (text or "").strip() or f"[No text detected by OCR] ({file_name})"
    return [text]


# Video loader (only audio used)
def extract_audio_from_video(file_name: str, data: bytes) -> bytes:
    """
    Extract mono 16 kHz WAV bytes from video using pydub+ffmpeg.
    """
    if not _HAVE_PYDUB:
        raise RuntimeError("pydub not installed. Run: pip install pydub (and ensure ffmpeg is on PATH)")

    fmt = {
        ".mp4": "mp4",
        ".mov": "mov",
        ".mkv": "matroska",
        ".webm": "webm",
    }.get(get_file_suffix(file_name), "mp4")

    # Load video bytes into pydub
    seg = AudioSegment.from_file(io.BytesIO(data), format=fmt)

    # Convert to mono 16kHz WAV PCM16
    seg = seg.set_frame_rate(16000).set_channels(1).set_sample_width(2)
    buf = io.BytesIO()

    # Export to WAV bytes
    seg.export(buf, format="wav")
    return buf.getvalue()


# ASR for audio
def transcribe_audio(data: bytes, file_name: str) -> str:
    """ Use ASR to transcribe audio to text. """

    # Use imported whisper_audio_to_text function from providers/asr.py
    if not _HAVE_ASR or whisper_audio_to_text is None:
        return "[ASR unavailable: local Whisper provider not importable]"
    
    # Load ASR config
    cfg = load_config() if _HAVE_CONFIG else {}
    a = (cfg.get("asr") or {})

    # Call the ASR function
    return whisper_audio_to_text(
    data=data,
    file_name=file_name,
    model=a.get("model", "tiny.en"),
    language=a.get("language"),
    device=a.get("device", "auto"),
    compute_type=a.get("compute_type", "int8"),
    vad=a.get("vad", True),
    beam_size=a.get("beam_size", 5),
    no_speech_threshold=a.get("no_speech_threshold", 0.6),
    reencode_fallback=a.get("reencode_fallback", True),
    )


# ---------------------- Main Loader function --------------------- #

def load_texts_from_bytes(file_name: str, data: bytes) -> Tuple[List[str], str]:
    """
    Returns:
      - texts: List[str] (later chunked in pipeline)
      - media_type: 'document' | 'presentation' | 'image' | 'audio' | 'video' | 'other'
    """

    # First get file extension - used to route to correct loader
    ext = get_file_suffix(file_name)

    # Plain text/markdown
    if ext in SUPPORTED_DOC:
        try:
            txt = data.decode("utf-8", errors="ignore")
            return [txt], "document"
        except Exception as e:
            return [f"[Text decode failed: {e}] ({file_name})"], "document"

    # PDF
    if ext in SUPPORTED_PDF:
        try:
            texts = load_text_from_pdf(file_name, data)
            if not texts:
                texts = [f"[Empty PDF or non-extractable text in {file_name}]"]
            return texts, "document"
        except Exception as e:
            return [f"[PDF extract error: {e}] ({file_name})"], "document"

    # Word
    if ext in SUPPORTED_WORD:
        if ext == ".docx":
            try:
                texts = load_text_from_docx(file_name, data)
                if not texts:
                    texts = [f"[Empty DOCX content] ({file_name})"]
                return texts, "document"
            except Exception as e:
                return [f"[DOCX extract error: {e}] ({file_name})"], "document"
        else:
            # legacy .doc
            return [f"[Legacy .doc not yet supported: convert to .docx] ({file_name})"], "document"

    # PowerPoint
    if ext in SUPPORTED_PRES:
        try:
            texts = load_text_from_ppt(file_name, data)
            return texts, "presentation"
        except Exception as e:
            return [f"[Presentation extract error: {e}] ({file_name})"], "presentation"

    # Images (OCR)
    if ext in SUPPORTED_IMAGE:
        try:
            return load_image_text(file_name, data), "image"
        except Exception as e:
            return [f"[Image OCR error for {file_name}: {e}]"], "image"

    # AUDIO → ASR
    if ext in SUPPORTED_AUDIO and ext not in SUPPORTED_VIDEO:
        try:
            text = transcribe_audio(data, file_name)
            return [text], "audio"
        except Exception as e:
            return [f"[Audio ASR error for {file_name}: {e}]"], "audio"

    # VIDEO → extract audio → ASR
    if ext in SUPPORTED_VIDEO:
        try:
            wav_bytes = extract_audio_from_video(file_name, data)
            text = transcribe_audio(wav_bytes, f"{Path(file_name).stem}.wav")
            return [text], "video"
        except Exception as e:
            return [f"[Video transcription error for {file_name}: {e}]"], "video"

    # Fallback
    return [f"[Unsupported file type placeholder for {file_name}]"], "other"
